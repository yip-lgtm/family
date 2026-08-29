#!/usr/bin/env python3
"""Generate a 3-minute pitch deck for 教父世家."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x0B, 0x08, 0x04)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
CREAM = RGBColor(0xF4, 0xE8, 0xC8)
MUTED = RGBColor(0xB8, 0xA8, 0x88)
CYAN = RGBColor(0x7E, 0xC8, 0xC3)
PANEL = RGBColor(0x14, 0x10, 0x0A)

OUT = Path(__file__).resolve().parents[1] / "docs" / "教父世家-3分鐘簡報.pptx"
TOTAL = 11


def set_run(run, text, size=18, bold=False, color=CREAM, font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = parse_xml(
        f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>'
    )
    rPr.append(ea)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def bg(slide, color=INK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def bar(slide, top, height=Inches(0.04), color=GOLD):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, W, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def box(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    else:
        sh.line.fill.background()
    return sh


def tb(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, bold, color) or str"""
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, 18, False, CREAM)
        text, size, bold, color = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        set_run(p.add_run(), text, size, bold, color)
    return tf


def seal(slide, l, t, size=Inches(0.86)):
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x16, 0x12, 0x0C)
    sh.line.color.rgb = GOLD
    sh.line.width = Pt(1.5)
    tb(slide, l, t + size * 0.28, size, size * 0.48, [("父", 22, True, GOLD)], PP_ALIGN.CENTER)
    return sh


def footer(slide, n):
    tb(slide, Inches(0.55), Inches(7.08), Inches(8.2), Inches(0.28),
       [("教父世家  ·  THE FAMILY", 11, False, MUTED)])
    tb(slide, Inches(10.4), Inches(7.08), Inches(2.3), Inches(0.28),
       [(f"{n:02d}  /  {TOTAL:02d}", 11, False, MUTED)], PP_ALIGN.RIGHT)


def new_slide(prs, n=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    bar(s, Inches(0), Inches(0.06))
    bar(s, H - Inches(0.06), Inches(0.06))
    if n is not None:
        footer(s, n)
    return s


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Title
    s = new_slide(prs)
    seal(s, Inches(6.23), Inches(0.85), Inches(0.88))
    tb(s, Inches(0.8), Inches(1.78), Inches(11.7), Inches(0.4),
       [("蒼梧山  ·  教父山門  ·  GODFATHER GATE", 15, False, GOLD)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(2.22), Inches(11.7), Inches(1.05),
       [("教父世家", 54, True, GOLD)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(3.22), Inches(11.7), Inches(0.4),
       [("THE FAMILY", 18, False, MUTED)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(3.72), Inches(11.7), Inches(0.6),
       [("你唔係點擊升級。你係天道。", 26, False, CREAM)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(1.7),
       [
           ("約 3 分鐘  ·  含製作 workflow  ·  演示畫先上", 15, False, MUTED),
           ("玩：https://yip-lgtm.github.io/family/", 18, True, CYAN),
           ("源碼：https://github.com/yip-lgtm/family.git", 16, False, MUTED),
       ], PP_ALIGN.CENTER)
    add_notes(s, """大家好。呢個係教父世家。品牌印係「父」，英文 THE FAMILY。地圖標題：教父山門，GODFATHER GATE。
你唔係點擊農夫。你入去係一間修仙世家，你係天道：睇住人，間中伸手改命。
而家就可以玩：yip-lgtm.github.io/family。今日三分鐘：做咩、畫面點樣永遠有人、點做出嚟。""")

    # 2 Not a clicker
    s = new_slide(prs, 2)
    tb(s, Inches(0.7), Inches(0.32), Inches(12), Inches(0.9),
       [("01  ·  唔係點擊遊戲", 14, False, GOLD),
        ("玩家角色：天道，唔係農夫", 32, True, CREAM)])
    cards = [
        ("傳統放置", "囤靈氣\n買升級\n睇數字漲", MUTED),
        ("教父世家", "養人\n改命\n睇戲", GOLD),
        ("你出手", "賜福 · 天劫\n心魔 · 指派\n耗的是氣運", CYAN),
    ]
    for i, (t, b, c) in enumerate(cards):
        x = Inches(0.7 + i * 4.1)
        box(s, x, Inches(1.85), Inches(3.8), Inches(4.7), RGBColor(0x16, 0x12, 0x0C), c)
        tb(s, x + Inches(0.25), Inches(2.1), Inches(3.3), Inches(0.6),
           [(t, 20, True, c)])
        tb(s, x + Inches(0.25), Inches(2.85), Inches(3.3), Inches(3.2),
           [(b, 22, False, CREAM)])
    add_notes(s, """傳統放置：你係農夫，囤靈氣、買升級、睇數字。
教父世家唔係。靈樞可以聚氣，但真正玩法係一班有根骨、心性、關係網嘅人。
你用氣運賜福、天劫、心魔、指派。氣運唔夠就唔好亂改。少出手，山門自己過月。你出手，就係改命。""")

    # 3 Heaven
    s = new_slide(prs, 3)
    tb(s, Inches(0.7), Inches(0.32), Inches(12), Inches(0.9),
       [("02  ·  天道介面", 14, False, GOLD),
        ("一屏睇晒山門、人、同你可改嘅命", 28, True, CREAM)])
    items = [
        ("曆法 HUD", "玄元曆、氣運、人口、最高境界。暫停、一倍、三倍、八倍過月。"),
        ("教父山門", "GODFATHER GATE。靈樞、後山、丹房、藏經閣、山門、雲市。人會自己揀去邊。"),
        ("人物檢視", "根骨、心性、功法、丹藥、法寶、心事、記憶、關係。開場畫像即見。"),
        ("天道四鍵", "賜福八氣運、天劫十二、心魔十、指派三。氣運唔夠就唔好亂改。"),
        ("家族傳承", "家訓、老祖突破、三選一天賦。金丹元嬰有金光。老祖唔會死。"),
        ("事件流", "功法、丹方、法寶、紀年。人死功法可以留下。"),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        x, y = Inches(0.55 + col * 4.2), Inches(1.55 + row * 2.55)
        box(s, x, y, Inches(4.0), Inches(2.35), PANEL, GOLD)
        tb(s, x + Inches(0.2), y + Inches(0.18), Inches(3.6), Inches(0.45), [(t, 18, True, GOLD)])
        tb(s, x + Inches(0.2), y + Inches(0.7), Inches(3.6), Inches(1.45), [(b, 15, False, CREAM)])
    add_notes(s, """介面係天道工作台。上面曆法、氣運、人口。時間可停，可一倍、三倍、八倍。
中間地圖標題：教父山門，GODFATHER GATE。人會自己去後山、丹房、藏經閣、雲市。
天道四鍵好貴：賜福八、天劫十二、心魔十、指派三。開場畫像即見。老祖唔會死。""")

    # 4 Living people
    s = new_slide(prs, 4)
    tb(s, Inches(0.7), Inches(0.26), Inches(12), Inches(0.85),
       [("03  ·  活人，唔係數值", 14, False, GOLD),
        ("開局六人 · 畫像已捆進遊戲", 28, True, CREAM)])
    people = [
        ("青玄機", "老祖 · 《青嵐吐納訣》", "鎮山門。你唔可以等佢死。"),
        ("沈清梧", "長老 · 穩、顧家", "會頂老祖，家要先立住。"),
        ("葉疏影", "二代 · 冷、準", "同清梧有舊怨，戲先有衝突。"),
        ("白無塵", "外來 · 身世成謎", "散修入譜，笑意難測。"),
        ("蒼小魚", "丹房 · 貪生識時務", "會偷雞，亦會留低嚟。"),
        ("嵐七七", "記名 · 細、靈", "成日問不該問嘅嘢。"),
    ]
    for i, (name, role, blurb) in enumerate(people):
        col, row = i % 3, i // 3
        x, y = Inches(0.55 + col * 4.2), Inches(1.28 + row * 2.7)
        box(s, x, y, Inches(4.0), Inches(2.5), PANEL, GOLD)
        tb(s, x + Inches(0.22), y + Inches(0.22), Inches(3.5), Inches(0.48), [(name, 22, True, GOLD)])
        tb(s, x + Inches(0.22), y + Inches(0.78), Inches(3.5), Inches(0.4), [(role, 14, False, CYAN)])
        tb(s, x + Inches(0.22), y + Inches(1.28), Inches(3.5), Inches(0.9), [(blurb, 16, False, CREAM)])
    add_notes(s, """開局六人，性格寫死，畫像已捆進遊戲，一開就見。
青玄機鎮山門，修《青嵐吐納訣》。沈清梧穩、顧家。葉疏影冷、同清梧有舊怨。白無塵外來。蒼小魚識時務。嵐七七成日問不該問嘅。
人會死。死咗，功法可以留低。""")

    # 5 Screenplay + AI
    s = new_slide(prs, 5)
    tb(s, Inches(0.7), Inches(0.26), Inches(12), Inches(0.85),
       [("04  ·  編劇組 + OpenRouter", 14, False, GOLD),
        ("《教父》三部曲節奏，唔抄對白", 28, True, CREAM)])
    acts = [
        ("第一幕", "血色開端", "權力真空、聯姻、暗湧。家要先立住。"),
        ("第二幕", "雙生歲月", "繼承、背叛、事業擴張。人開始唔同路。"),
        ("第三幕", "最後輓歌", "贖罪、遺產、落幕。唔一定團圓。"),
    ]
    for i, (act, title, blurb) in enumerate(acts):
        x = Inches(0.55 + i * 4.2)
        box(s, x, Inches(1.28), Inches(4.0), Inches(2.4), PANEL, GOLD)
        tb(s, x + Inches(0.2), Inches(1.42), Inches(3.6), Inches(0.32), [(act, 13, False, MUTED)])
        tb(s, x + Inches(0.2), Inches(1.78), Inches(3.6), Inches(0.45), [(title, 22, True, GOLD)])
        tb(s, x + Inches(0.2), Inches(2.32), Inches(3.6), Inches(1.1), [(blurb, 15, False, CREAM)])
    box(s, Inches(0.55), Inches(3.9), Inches(12.2), Inches(2.85), PANEL, CYAN)
    tb(s, Inches(0.8), Inches(4.08), Inches(11.7), Inches(2.5), [
        ("現場：每兩個月自動一場戲 · 可按「下一場」催更", 16, True, CYAN),
        ("有 Key：文字預設 openrouter/free。編劇喺背景無限連載，冷卻更短。Key 只喺瀏覽器。", 15, False, CREAM),
        ("冇 Key / 失敗：劇組代班——仍然係完整劇本，唔會停機。", 15, False, CREAM),
        ("致敬結構同主題，唔抄電影對白，唔抄電影人名。戲會寫入心事同關係。", 15, False, CREAM),
    ])
    add_notes(s, """編劇組結構致敬《教父》三部曲：血色開端、雙生歲月、最後輓歌。致敬節奏，唔抄對白、唔抄人名。
每兩個月自動一場戲。文字預設 openrouter/free，Key 只存瀏覽器。有 Key，編劇背景無限連載，冷卻更短。
冇 Key：劇組代班，戲照拍，寫入心事同關係。""")

    # 6 Demo art + LLM illustrator  (NEW)
    s = new_slide(prs, 6)
    tb(s, Inches(0.7), Inches(0.26), Inches(12), Inches(0.85),
       [("05  ·  畫面永遠有人", 14, False, GOLD),
        ("演示畫先上，LLM 畫好先換", 28, True, CREAM)])
    arts = [
        ("開場六人畫像", "青玄機、沈清梧、葉疏影、白無塵、蒼小魚、嵐七七。捆進遊戲，一開即見。"),
        ("場次演示畫", "大殿夜宴、雨中石階、山門、燈密室。場戲一出就有畫面。"),
        ("背景無限出圖", "有 Key：編劇同畫師喺背景不停寫、不停畫。512 檔細圖，唔擋操作。"),
        ("靜默重試", "失敗就自己再試。演示畫留住，直到 LLM 畫換上去。唔會空白「插畫未成」。"),
    ]
    for i, (t, b) in enumerate(arts):
        col, row = i % 2, i // 2
        x, y = Inches(0.55 + col * 6.35), Inches(1.28 + row * 2.15)
        box(s, x, y, Inches(6.1), Inches(2.0), PANEL, GOLD)
        tb(s, x + Inches(0.28), y + Inches(0.2), Inches(5.55), Inches(0.45), [(t, 18, True, GOLD)])
        tb(s, x + Inches(0.28), y + Inches(0.72), Inches(5.55), Inches(1.1), [(b, 15, False, CREAM)])
    box(s, Inches(0.55), Inches(5.68), Inches(12.2), Inches(1.12), PANEL, CYAN)
    tb(s, Inches(0.8), Inches(5.82), Inches(11.7), Inches(0.9), [
        ("插畫預設 google/gemini-2.5-flash-image  ·  模型設定可關「自動生成插畫」  ·  冇 Key 就留演示畫", 15, False, CYAN),
    ])
    add_notes(s, """畫面最怕空白「插畫未成」。所以六人開場畫像同場次演示畫，全部捆進遊戲。
有 Key：編劇同畫師喺背景無限開工。畫師預設 google/gemini-2.5-flash-image，512 檔細圖，失敗靜默重試。
演示畫留住，LLM 畫好先換。模型設定可以關自動出圖。冇 Key 或者失敗：演示畫唔走。你永遠唔會睇住一個空框。""")

    # 7 Craft
    s = new_slide(prs, 7)
    tb(s, Inches(0.7), Inches(0.26), Inches(12), Inches(0.85),
       [("06  ·  金漆山門，現場聲畫", 14, False, GOLD),
        ("CSS 牌坊 · Web Audio · 嗶哩官方嵌入", 26, True, CREAM)])
    craft = [
        ("視覺", "CSS 金漆牌坊、墨底描金、品牌印「父」。Canvas 靈塵。突破金光。演示畫打進倉庫，唔靠外鏈圖床。"),
        ("音效", "瀏覽器合成木魚、磬、雷、心魔。一鍵靜音。唔使下載音檔。"),
        ("BGM", "嗶哩嗶哩官方嵌入《關注塔菲喵》循環歌單。可縮小成一條。唔下載、唔轉檔、唔盜鏈音源。"),
        ("技術", "Vite + 原生 JS。規則、編劇、畫師分檔。GitHub Pages 可直接開。"),
    ]
    for i, (t, b) in enumerate(craft):
        col, row = i % 2, i // 2
        x, y = Inches(0.55 + col * 6.35), Inches(1.28 + row * 2.6)
        box(s, x, y, Inches(6.1), Inches(2.4), PANEL, GOLD)
        tb(s, x + Inches(0.3), y + Inches(0.25), Inches(5.5), Inches(0.45), [(t, 22, True, GOLD)])
        tb(s, x + Inches(0.3), y + Inches(0.85), Inches(5.5), Inches(1.3), [(b, 16, False, CREAM)])
    add_notes(s, """牌坊係 CSS 描金，印係「父」。演示畫打進倉庫，唔靠外鏈圖床。音效 Web Audio 合成木魚磬雷。
BGM 係嗶哩官方嵌入，唔下載、唔轉檔、唔盜鏈音源。技術棧：Vite 加原生 JS。""")

    # 8 Production workflow
    s = new_slide(prs, 8)
    tb(s, Inches(0.7), Inches(0.24), Inches(12), Inches(0.85),
       [("07  ·  教父世家製作 workflow", 14, False, GOLD),
        ("先立天道，再寫規則，畫面先演示再 LLM", 24, True, CREAM)])
    flow = [
        ("1", "立意", "玩家係天道。\n養人、改命、睇戲。\n唔做點擊農夫。"),
        ("2", "規則", "world.js\n月結、根骨、壽元\n氣運、死亡、傳承。"),
        ("3", "編劇", "screenplay.js\n三幕節奏 + 劇組代班。\n有 Key 就背景無限連載。"),
        ("4", "皮與聲", "先捆六人+場次演示畫。\n再接 LLM 背景出圖。\nCSS、Web Audio、嗶哩嵌入。"),
        ("5", "試玩", "Vite 開 /family/\n改完即睇。\n細畫面、BGM、文案。"),
        ("6", "出貨", "npm run build\n推 gh-pages。\ngithub.io/family"),
    ]
    for i, (n, t, b) in enumerate(flow):
        col, row = i % 3, i // 3
        x, y = Inches(0.5 + col * 4.2), Inches(1.22 + row * 2.7)
        box(s, x, y, Inches(4.0), Inches(2.5), PANEL, GOLD)
        tb(s, x + Inches(0.2), y + Inches(0.16), Inches(3.6), Inches(0.4),
           [(n + "  ·  " + t, 18, True, GOLD)])
        tb(s, x + Inches(0.2), y + Inches(0.65), Inches(3.6), Inches(1.65),
           [(line, 15, False, CREAM) for line in b.split("\n")])
    add_notes(s, """製作唔係先畫皮。六步：立意，玩家係天道。規則 world.js。編劇 screenplay.js，三幕加劇組代班；有 Key 就背景無限連載。
皮與聲：先捆六人同場次演示畫，再接 LLM 背景出圖。試玩 Vite 開 /family/。出貨 build 推 gh-pages，就係而家呢個站。""")

    # 9 Engineering loop
    s = new_slide(prs, 9)
    tb(s, Inches(0.7), Inches(0.24), Inches(12), Inches(0.85),
       [("08  ·  工程回路", 14, False, GOLD),
        ("六個檔，一條每日回路：改 → 睇 → 推", 26, True, CREAM)])
    modules = [
        ("world.js", "月結、人物、天道四鍵"),
        ("screenplay.js", "導演、三幕、OpenRouter"),
        ("main.js", "HUD、地圖、檢視、操作"),
        ("style.css", "金漆牌坊、墨底描金、父印"),
        ("fx.js", "靈塵、聚氣、突破金光"),
        ("audio.js", "木魚磬雷 + 嗶哩嵌入"),
    ]
    for i, (name, blurb) in enumerate(modules):
        col, row = i % 3, i // 3
        x, y = Inches(0.5 + col * 4.2), Inches(1.2 + row * 1.58)
        box(s, x, y, Inches(4.0), Inches(1.42), PANEL, GOLD)
        tb(s, x + Inches(0.2), y + Inches(0.16), Inches(3.6), Inches(0.38),
           [(name, 18, True, CYAN)])
        tb(s, x + Inches(0.2), y + Inches(0.62), Inches(3.6), Inches(0.55),
           [(blurb, 15, False, CREAM)])
    box(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.28), PANEL, CYAN)
    tb(s, Inches(0.75), Inches(4.64), Inches(11.8), Inches(2.0), [
        ("畫面　demoArt.js 捆演示畫 → illustrate.js 背景出 512 檔圖 → 靜默重試 → 畫好先換", 16, True, CYAN),
        ("日常　改規則或文案 → npm run dev → 現場睇一場戲 → git commit", 15, False, CREAM),
        ("出貨　npm run build → 將 dist 推去 gh-pages → https://yip-lgtm.github.io/family/", 15, False, CREAM),
        ("約束　唔下載嗶哩音源 · 唔抄《教父》對白 · Key 只存瀏覽器 · 演示畫永遠兜底", 14, False, MUTED),
    ])
    add_notes(s, """核心六個檔。world 規則，screenplay 戲，main 介面，style fx audio 現場。
畫面掛兩件：demoArt 捆演示畫，illustrate 背景出 512 檔圖，失敗靜默重試。
紅線：唔偷嗶哩音、唔抄教父對白、Key 唔上伺服器。演示畫永遠兜底。""")

    # 10 How to play 3 min
    s = new_slide(prs, 10)
    tb(s, Inches(0.7), Inches(0.26), Inches(12), Inches(0.85),
       [("09  ·  而家點玩（三分鐘）", 14, False, GOLD),
        ("四步入局，唔使註冊", 28, True, CREAM)])
    steps = [
        ("1", "打開連結", "yip-lgtm.github.io/family\n撳「承天命」。"),
        ("2", "睇一個人", "點沈清梧或者葉疏影。\n畫像即見，記住心事。"),
        ("3", "等兩場戲", "開時間。睇編劇組出第一、第二場——場次都有演示畫。"),
        ("4", "先改一次命", "賜福或者天劫。\n然後停手，睇後果。"),
    ]
    for i, (n, t, b) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        box(s, x, Inches(1.32), Inches(3.0), Inches(4.85), PANEL, GOLD)
        tb(s, x + Inches(0.2), Inches(1.55), Inches(2.6), Inches(0.7), [(n, 36, True, GOLD)], PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.65), [(t, 18, True, CREAM)], PP_ALIGN.CENTER)
        tb(s, x + Inches(0.15), Inches(3.15), Inches(2.7), Inches(2.6), [(b, 15, False, MUTED)], PP_ALIGN.CENTER)
    add_notes(s, """四步，唔使註冊。打開 yip-lgtm.github.io/family，撳承天命。點沈清梧或者葉疏影——畫像即見。
開時間，等編劇組兩場。只用一次天道：賜福或者天劫，然後停手。
可選貼 OpenRouter Key，編劇同畫師就喺背景開工。冇 Key 一樣完整。""")

    # 11 Close
    s = new_slide(prs)
    seal(s, Inches(6.23), Inches(0.72), Inches(0.82))
    tb(s, Inches(0.8), Inches(1.62), Inches(11.7), Inches(0.4),
       [("結  ·  家訓", 14, False, GOLD)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(2.08), Inches(11.7), Inches(1.25),
       [("人死功法在。家在，戲就未完。", 32, True, GOLD)], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(3.45), Inches(11.7), Inches(1.7),
       [
           ("玩　　https://yip-lgtm.github.io/family/", 20, True, CYAN),
           ("源碼　https://github.com/yip-lgtm/family.git", 18, False, MUTED),
           ("鏡像　https://github.com/yip-lgtm/qi-lineage", 16, False, MUTED),
       ], PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.0),
       [
           ("開源 · 瀏覽器即玩 · Key 可選 · 演示畫永遠在 · 劇組永遠代班", 16, False, MUTED),
           ("THE FAMILY  ·  教父山門  ·  GODFATHER GATE", 14, False, GOLD),
       ], PP_ALIGN.CENTER)
    add_notes(s, """收束。放置遊戲可以係一場有人、有怨、有遺產嘅戲。你係天道。少出手。睇住人。畫面永遠有人。
人死功法在。家在，戲就未完。github.io/family。多謝。""")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
